import io

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from reviewer.parsing.base import OcrFn, ParsedContent


def _shape_text(shape, ocr: OcrFn, parts: list[str]) -> None:
    if shape.has_text_frame and shape.text_frame.text.strip():
        parts.append(shape.text_frame.text.strip())
    if shape.has_table:
        for row in shape.table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append("\t".join(cells))
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            image = shape.image
            transcription = ocr(image.blob, image.content_type)
        except Exception:
            transcription = ""
        if transcription.strip():
            parts.append(transcription.strip())


def read_pptx(data: bytes, ocr: OcrFn) -> ParsedContent:
    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            _shape_text(shape, ocr, parts)
        if slide.has_notes_slide:
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                parts.append(f"[Notes] {note}")
    return ParsedContent(text="\n\n".join(parts))
