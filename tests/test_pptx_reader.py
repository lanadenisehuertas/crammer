import io
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
from reviewer.parsing.pptx_reader import read_pptx


def _png() -> io.BytesIO:
    buf = io.BytesIO()
    Image.new("RGB", (8, 8), "white").save(buf, format="PNG")
    buf.seek(0)
    return buf


def _make_pptx() -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Slide Title"
    slide.shapes.add_picture(_png(), Inches(1), Inches(1))
    slide.notes_slide.notes_text_frame.text = "Speaker note here"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_read_pptx_extracts_text_notes_and_image_ocr():
    pc = read_pptx(_make_pptx(), ocr=lambda b, m: "SLIDE_IMAGE_TEXT")
    assert "Slide Title" in pc.text
    assert "Speaker note here" in pc.text
    assert "SLIDE_IMAGE_TEXT" in pc.text
